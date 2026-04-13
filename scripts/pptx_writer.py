"""
PPTX 번역 텍스트 삽입 — 원본 서식 보존.

PPT 워크플로 Step 4: 원본 PPTX를 열어 slides_manifest.json의 구조와
대응되는 번역 텍스트를 run 단위로 삽입한다.
run.font 속성(크기/굵기/이탤릭/색상)은 유지하고 텍스트만 교체한다.
영문 폰트는 한글 대응 폰트로 교체한다.
"""

from __future__ import annotations

import argparse
import copy
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print(
        "Install dependencies: pip install python-pptx  "
        "(or pip install -r scripts/requirements.txt)",
        file=sys.stderr,
    )
    raise


FONT_MAP: dict[str | None, str] = {
    "Arial": "Malgun Gothic",
    "Calibri": "Malgun Gothic",
    "Times New Roman": "Malgun Gothic",
    "Helvetica": "Malgun Gothic",
    "Verdana": "Malgun Gothic",
    "Cambria": "Malgun Gothic",
    "Segoe UI": "Malgun Gothic",
    None: "Malgun Gothic",
}


def project_root() -> Path:
    return Path(__file__).resolve().parent.parent


def load_translations(path: Path) -> dict:
    """번역 JSON 로드. 구조: { slide_number(str): { shape_id(str): [번역된 para 텍스트 리스트] } }"""
    return json.loads(path.read_text(encoding="utf-8"))


def map_font(original: str | None) -> str:
    """원본 영문 폰트를 한글 대응 폰트로 매핑."""
    if original and any(ord(c) > 0xAC00 for c in original):
        return original
    return FONT_MAP.get(original, FONT_MAP[None])


def apply_translations(
    pptx_path: Path,
    translations: dict,
    out_path: Path,
    *,
    remap_fonts: bool = True,
) -> tuple[int, int]:
    """
    원본 PPTX에 번역 텍스트를 삽입하고 저장한다.

    translations 구조:
    {
      "1": {                  # slide_number (1-based, 문자열)
        "42": [               # shape_id (문자열)
          "첫 번째 단락 번역",
          "두 번째 단락 번역"
        ]
      }
    }

    반환: (수정된 슬라이드 수, 수정된 도형 수)
    """
    prs = Presentation(str(pptx_path))
    modified_slides = 0
    modified_shapes = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_key = str(slide_idx)
        if slide_key not in translations:
            continue

        slide_trans = translations[slide_key]
        slide_touched = False

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            shape_key = str(shape.shape_id)
            if shape_key not in slide_trans:
                continue

            translated_paras: list[str] = slide_trans[shape_key]
            tf = shape.text_frame

            for para_idx, para in enumerate(tf.paragraphs):
                if para_idx >= len(translated_paras):
                    break

                new_text = translated_paras[para_idx]
                if not para.runs:
                    continue

                if len(para.runs) == 1:
                    run = para.runs[0]
                    run.text = new_text
                    if remap_fonts:
                        run.font.name = map_font(run.font.name)
                else:
                    para.runs[0].text = new_text
                    if remap_fonts:
                        para.runs[0].font.name = map_font(para.runs[0].font.name)
                    for extra_run in para.runs[1:]:
                        extra_run.text = ""

            modified_shapes += 1
            slide_touched = True

        if slide_touched:
            modified_slides += 1

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return modified_slides, modified_shapes


def main() -> int:
    p = argparse.ArgumentParser(
        description="Insert translated text into PPTX (preserving formatting)."
    )
    p.add_argument("pptx", type=Path, help="Original .pptx file")
    p.add_argument("translations", type=Path, help="Translations JSON file")
    p.add_argument(
        "-o", "--output", type=Path, default=None,
        help="Output .pptx path (default: output/{slug}/final_translated.pptx)",
    )
    p.add_argument(
        "--no-font-remap", action="store_true",
        help="Skip remapping English fonts to Korean fonts",
    )
    args = p.parse_args()

    pptx_path: Path = args.pptx.resolve()
    trans_path: Path = args.translations.resolve()

    if not pptx_path.exists():
        print(f"PPTX not found: {pptx_path}", file=sys.stderr)
        return 1
    if not trans_path.exists():
        print(f"Translations file not found: {trans_path}", file=sys.stderr)
        return 1

    root = project_root()
    slug = pptx_path.stem
    out_path = args.output or (root / "output" / slug / "final_translated.pptx")

    translations = load_translations(trans_path)
    slides, shapes = apply_translations(
        pptx_path, translations, out_path,
        remap_fonts=not args.no_font_remap,
    )
    print(f"Done: {slides} slides, {shapes} shapes modified → {out_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
