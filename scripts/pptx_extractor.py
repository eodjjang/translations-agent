"""
PPTX 슬라이드 텍스트 추출 → slides_manifest.json 생성.

PPT 워크플로 Step 1: 각 슬라이드의 텍스트 프레임을 순회하여
슬라이드 번호·도형 ID·원문 텍스트·폰트 정보를 JSON으로 저장한다.
이미지/도표/차트 내부 텍스트는 대상 아님.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print(
        "Install dependencies: pip install python-pptx  "
        "(or pip install -r scripts/requirements.txt)",
        file=sys.stderr,
    )
    raise


def project_root() -> Path:
    return Path(__file__).resolve().parent.parent


def slug_from_stem(stem: str) -> str:
    return stem


def _run_color_rgb(font) -> str | None:
    """Theme/scheme colors may not expose .rgb; return None if unavailable."""
    try:
        c = font.color
        if c is None:
            return None
        rgb = c.rgb
        return str(rgb) if rgb is not None else None
    except (AttributeError, TypeError, ValueError):
        return None


def extract_slides(pptx_path: Path) -> list[dict]:
    """슬라이드별 텍스트 프레임 정보를 추출한다."""
    prs = Presentation(str(pptx_path))
    slides_data: list[dict] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        shapes_data: list[dict] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            paragraphs: list[dict] = []
            for para in shape.text_frame.paragraphs:
                runs: list[dict] = []
                for run in para.runs:
                    font = run.font
                    run_info: dict = {
                        "text": run.text,
                        "font_name": font.name,
                        "font_size_pt": round(font.size.pt, 1) if font.size else None,
                        "bold": font.bold,
                        "italic": font.italic,
                    }
                    crgb = _run_color_rgb(font)
                    if crgb:
                        run_info["color_rgb"] = crgb
                    runs.append(run_info)

                paragraphs.append({
                    "text": para.text,
                    "level": para.level,
                    "runs": runs,
                })

            full_text = shape.text_frame.text.strip()
            if not full_text:
                continue

            shapes_data.append({
                "shape_id": shape.shape_id,
                "shape_name": shape.name,
                "text": full_text,
                "paragraphs": paragraphs,
            })

        slides_data.append({
            "slide_number": slide_idx,
            "shapes": shapes_data,
        })

    return slides_data


def write_manifest(slides_data: list[dict], out_path: Path) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    manifest = {
        "total_slides": len(slides_data),
        "slides": slides_data,
    }
    out_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")


def main() -> int:
    p = argparse.ArgumentParser(description="Extract text from PPTX slides.")
    p.add_argument("pptx", type=Path, help="Path to .pptx file")
    p.add_argument(
        "-o", "--output-dir", type=Path, default=None,
        help="Output directory (default: output/{slug}/)",
    )
    args = p.parse_args()

    pptx_path: Path = args.pptx.resolve()
    if not pptx_path.exists():
        print(f"File not found: {pptx_path}", file=sys.stderr)
        return 1

    slug = slug_from_stem(pptx_path.stem)
    root = project_root()
    out_dir = args.output_dir or (root / "output" / slug)
    out_dir.mkdir(parents=True, exist_ok=True)

    print(f"Extracting text from: {pptx_path.name}")
    slides_data = extract_slides(pptx_path)

    manifest_path = out_dir / "slides_manifest.json"
    write_manifest(slides_data, manifest_path)

    total_shapes = sum(len(s["shapes"]) for s in slides_data)
    print(f"Done: {len(slides_data)} slides, {total_shapes} text shapes → {manifest_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
