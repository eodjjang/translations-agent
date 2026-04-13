"""
PPT 서식 검증기: final_translated.pptx 저장 시 서식 보존 여부를 자동 검증한다.

검증 항목:
1. 텍스트 프레임 수 대응 — 원본 vs 번역본의 슬라이드별 텍스트 프레임 수 비교
2. 빈 텍스트 프레임 — 번역 후 비어있는 프레임이 없는지 확인
3. 폰트 보존 — run 단위 font.size, font.bold, font.italic 원본 동일 여부
4. References 슬라이드 미번역 — References 슬라이드 텍스트가 원본과 동일한지 확인
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    # python-pptx가 없으면 검증을 건너뛴다 (hook이 깨지지 않도록)
    Presentation = None


def _find_project_root(pptx_path: Path) -> Path:
    p = pptx_path.resolve()
    for parent in p.parents:
        if (parent / "input").is_dir() and (parent / "output").is_dir():
            return parent
    return p.parent.parent.parent


def _find_original_pptx(root: Path) -> Path | None:
    """input/ 폴더에서 원본 .pptx를 찾는다."""
    input_dir = root / "input"
    if not input_dir.is_dir():
        return None
    pptx_files = list(input_dir.glob("*.pptx"))
    return pptx_files[0] if pptx_files else None


def _get_slide_frames(prs) -> list[list[dict]]:
    """슬라이드별 텍스트 프레임 정보를 수집한다."""
    result = []
    for slide in prs.slides:
        frames = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            runs_info = []
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    runs_info.append({
                        "text": run.text,
                        "size": run.font.size,
                        "bold": run.font.bold,
                        "italic": run.font.italic,
                    })
            frames.append({
                "shape_id": shape.shape_id,
                "text": shape.text_frame.text.strip(),
                "runs": runs_info,
            })
        result.append(frames)
    return result


def _is_references_slide(frames: list[dict]) -> bool:
    """슬라이드가 References인지 판별한다."""
    for f in frames:
        text_lower = f["text"].lower()
        if re.match(r"^\s*references?\s*$", text_lower):
            return True
        if "references" in text_lower and len(f["text"]) < 30:
            return True
    return False


def validate(file_path: str, hook_input: dict) -> list[str]:
    """
    번역된 .pptx 파일을 원본과 비교 검증하고 오류 목록을 반환한다.
    오류가 없으면 빈 리스트를 반환한다.
    """
    if Presentation is None:
        return ["PPT 검증 건너뜀 — python-pptx가 설치되지 않았습니다."]

    errors: list[str] = []
    translated_path = Path(file_path).resolve()

    if not translated_path.exists():
        return errors

    root = _find_project_root(translated_path)
    original_path = _find_original_pptx(root)

    if not original_path:
        errors.append("PPT 검증 — input/ 폴더에서 원본 PPTX를 찾을 수 없습니다.")
        return errors

    try:
        orig_prs = Presentation(str(original_path))
        trans_prs = Presentation(str(translated_path))
    except Exception as e:
        errors.append(f"PPT 검증 — 파일 열기 실패: {e}")
        return errors

    orig_frames = _get_slide_frames(orig_prs)
    trans_frames = _get_slide_frames(trans_prs)

    # --- 1. 슬라이드 수 비교 ---
    if len(orig_frames) != len(trans_frames):
        errors.append(
            f"슬라이드 수 불일치 — 원본: {len(orig_frames)}장, "
            f"번역본: {len(trans_frames)}장"
        )
        return errors

    frame_count_mismatches = []
    empty_frames = []
    font_mismatches = []
    ref_modified = []

    for slide_idx, (orig_slide, trans_slide) in enumerate(
        zip(orig_frames, trans_frames), start=1
    ):
        # --- 2. 텍스트 프레임 수 ---
        if len(orig_slide) != len(trans_slide):
            frame_count_mismatches.append(
                f"슬라이드 {slide_idx}: 원본 {len(orig_slide)}개 vs "
                f"번역본 {len(trans_slide)}개"
            )

        # --- 3. 빈 텍스트 프레임 ---
        for fi, frame in enumerate(trans_slide):
            if not frame["text"] and fi < len(orig_slide) and orig_slide[fi]["text"]:
                empty_frames.append(
                    f"슬라이드 {slide_idx}, 프레임 {fi + 1}: "
                    f"원본에 텍스트가 있었으나 번역본이 비어 있음"
                )

        # --- 4. 폰트 보존 ---
        for fi, (of, tf) in enumerate(
            zip(orig_slide, trans_slide[: len(orig_slide)])
        ):
            for ri, (orun, trun) in enumerate(
                zip(of["runs"], tf["runs"][: len(of["runs"])])
            ):
                issues = []
                if orun["size"] != trun["size"] and orun["size"] is not None:
                    issues.append("크기")
                if orun["bold"] != trun["bold"] and orun["bold"] is not None:
                    issues.append("굵기")
                if orun["italic"] != trun["italic"] and orun["italic"] is not None:
                    issues.append("이탤릭")
                if issues:
                    font_mismatches.append(
                        f"슬라이드 {slide_idx}, 프레임 {fi + 1}, "
                        f"run {ri + 1}: {', '.join(issues)} 변경됨"
                    )

        # --- 5. References 미번역 ---
        if _is_references_slide(orig_slide):
            for of, tf in zip(orig_slide, trans_slide[: len(orig_slide)]):
                if of["text"] != tf["text"]:
                    ref_modified.append(f"슬라이드 {slide_idx}")
                    break

    if frame_count_mismatches:
        sample = frame_count_mismatches[:3]
        more = (
            f" 외 {len(frame_count_mismatches) - 3}건"
            if len(frame_count_mismatches) > 3
            else ""
        )
        errors.append(
            f"텍스트 프레임 수 불일치 — {len(frame_count_mismatches)}건: "
            + "; ".join(sample) + more
        )

    if empty_frames:
        sample = empty_frames[:3]
        more = (
            f" 외 {len(empty_frames) - 3}건" if len(empty_frames) > 3 else ""
        )
        errors.append(
            f"빈 텍스트 프레임 — {len(empty_frames)}건: "
            + "; ".join(sample) + more
        )

    if font_mismatches:
        sample = font_mismatches[:3]
        more = (
            f" 외 {len(font_mismatches) - 3}건" if len(font_mismatches) > 3 else ""
        )
        errors.append(
            f"폰트 서식 변경 — {len(font_mismatches)}건: "
            + "; ".join(sample) + more
        )

    if ref_modified:
        errors.append(
            f"References 슬라이드 변경됨 — 원문 유지되어야 합니다: "
            + ", ".join(ref_modified)
        )

    return errors
